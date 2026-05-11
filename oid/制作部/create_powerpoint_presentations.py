#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint形式のプレゼン資料を自動生成（絵柄・カラーデザイン付き）
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class PresentationDesigner:
    def __init__(self):
        base_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        self.base_path = base_path

        # カラースキーム
        self.colors = {
            'primary': RGBColor(102, 126, 234),      # 紫青
            'secondary': RGBColor(118, 75, 162),     # 深紫
            'accent': RGBColor(255, 152, 0),         # オレンジ
            'text': RGBColor(44, 62, 80),            # ダークグレー
            'light_bg': RGBColor(245, 245, 245),     # 薄いグレー
            'white': RGBColor(255, 255, 255)
        }

    def create_sales_presentation(self):
        """AI営業自動化 - PowerPoint資料"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # スライド1: タイトル
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 白紙レイアウト
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['primary']

        # タイトルテキスト
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = "AI時代の個人スキル販売術"
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER

        # サブタイトル
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "ChatGPT/Geminiで営業時間を90%削減"
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER

        # スライド2: 問題定義
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['white']

        # ヘッダー
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        header_frame = header_box.text_frame
        header_frame.text = "従来型営業の課題"
        p = header_frame.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        # コンテンツボックス1
        content_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(4), Inches(5.5))
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = self.colors['light_bg']
        content_box.line.color.rgb = self.colors['primary']
        content_box.line.width = Pt(2)

        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)

        content = "営業担当者の1日:\n\n• メール作成: 2-3時間\n• リスト整理: 1時間\n• 返信対応: 1-2時間\n• フォロー: 1-2時間\n\n有意義な営業: 1時間"
        text_frame.text = content
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = self.colors['text']

        # コンテンツボックス2
        content_box2 = slide.shapes.add_shape(1, Inches(5.5), Inches(1.3), Inches(4), Inches(5.5))
        content_box2.fill.solid()
        content_box2.fill.fore_color.rgb = RGBColor(144, 238, 144)
        content_box2.line.color.rgb = self.colors['secondary']
        content_box2.line.width = Pt(2)

        text_frame2 = content_box2.text_frame
        text_frame2.word_wrap = True
        text_frame2.margin_top = Inches(0.2)
        text_frame2.margin_left = Inches(0.2)
        text_frame2.margin_right = Inches(0.2)

        content2 = "AI営業の効果:\n\n• メール生成: 3分\n• 返信分析: 1分\n• 月間接触: 50→500社\n• 成約率: 2%→5%\n\n有意義な営業: 6時間"
        text_frame2.text = content2
        for paragraph in text_frame2.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = paragraph == text_frame2.paragraphs[0]
            paragraph.font.color.rgb = self.colors['text']

        # スライド3: 30日ロードマップ
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['white']

        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        header_frame = header_box.text_frame
        header_frame.text = "30日間のロードマップ"
        p = header_frame.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']

        # ロードマップアイテム
        phases = [
            ("Day 1-2", "基礎習得", "ChatGPT/Gemini基本操作"),
            ("Day 3-7", "最適化", "プロンプト精度向上"),
            ("Day 8-14", "実運用", "毎日メール配信"),
            ("Day 15-21", "分析・改善", "返信率向上"),
            ("Day 22-28", "自動化完成", "ルーチン化"),
            ("Day 29-30", "評価・計画", "月間評価と改善")
        ]

        y_pos = 1.3
        for i, (day, phase, content) in enumerate(phases):
            # 背景
            box = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.9))
            box.fill.solid()
            box.fill.fore_color.rgb = self.colors['light_bg'] if i % 2 == 0 else self.colors['white']
            box.line.color.rgb = self.colors['primary']
            box.line.width = Pt(1)

            # テキスト
            text_frame = box.text_frame
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_top = Inches(0.1)

            p = text_frame.paragraphs[0]
            p.text = f"{day}  |  {phase}  |  {content}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['text']

            y_pos += 1.0

        return prs

    def create_sns_presentation(self):
        """SNS自動化 - PowerPoint資料"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # タイトルスライド
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['secondary']

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = "SNS運用自動化キット"
        p = title_frame.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "毎日5分で投稿完成、フォロワー月100人増加"
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER

        # 効果比較スライド
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['white']

        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        header_frame = header_box.text_frame
        header_frame.text = "AI投稿自動化の効果"
        p = header_frame.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.colors['secondary']

        # メトリクス表
        metrics = [
            ("投稿作成時間", "30分→5分", "83%削減"),
            ("月間投稿数", "10→30", "3倍"),
            ("フォロワー増加", "月5人→月100人", "20倍"),
            ("エンゲージ率", "2%→6%", "3倍"),
        ]

        y_pos = 1.4
        for metric, result, rate in metrics:
            box = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(1.2))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(230, 240, 255)
            box.line.color.rgb = self.colors['secondary']
            box.line.width = Pt(2)

            text_frame = box.text_frame
            text_frame.margin_left = Inches(0.3)
            text_frame.margin_top = Inches(0.15)

            p = text_frame.paragraphs[0]
            p.text = f"📊 {metric}"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.colors['secondary']

            p = text_frame.add_paragraph()
            p.text = f"  {result}  →  {rate}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['accent']
            p.level = 0

            y_pos += 1.3

        # 曜日別戦略スライド
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['white']

        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        header_frame = header_box.text_frame
        header_frame.text = "曜日別最適投稿パターン"
        p = header_frame.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.colors['secondary']

        days_content = [
            ("月", "トレンド情報"),
            ("火", "How-Toガイド"),
            ("水", "ビジュアル投稿"),
            ("木", "ストーリー型"),
            ("金", "励まし・モチベーション"),
            ("土", "知識・教育"),
            ("日", "目標設定"),
        ]

        x_pos = 0.5
        y_pos = 1.4
        col_count = 0

        for day, content in days_content:
            box = slide.shapes.add_shape(1, Inches(x_pos), Inches(y_pos), Inches(1.2), Inches(5))
            box.fill.solid()
            box.fill.fore_color.rgb = [
                RGBColor(200, 220, 255),
                RGBColor(220, 200, 255),
                RGBColor(255, 220, 200),
                RGBColor(200, 255, 200),
                RGBColor(255, 255, 200),
                RGBColor(255, 200, 220),
                RGBColor(220, 255, 255)
            ][col_count % 7]
            box.line.color.rgb = self.colors['primary']
            box.line.width = Pt(2)

            text_frame = box.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.3)

            p = text_frame.paragraphs[0]
            p.text = day
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self.colors['text']
            p.alignment = PP_ALIGN.CENTER

            p = text_frame.add_paragraph()
            p.text = content
            p.font.size = Pt(11)
            p.font.color.rgb = self.colors['text']
            p.alignment = PP_ALIGN.CENTER

            x_pos += 1.35
            col_count += 1

        return prs

    def create_ai_guide_presentation(self):
        """初心者向けAI活用ガイド - PowerPoint資料"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # タイトルスライド
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(76, 175, 80)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = "初心者向けAI活用ガイド"
        p = title_frame.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "ChatGPT/Gemini完全実装マニュアル"
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER

        # 学習ロードマップスライド
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors['white']

        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        header_frame = header_box.text_frame
        header_frame.text = "30日間の学習ロードマップ"
        p = header_frame.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(76, 175, 80)

        roadmap = [
            ("Day 1-5", "ChatGPT基本操作", "登録・ログイン・初回体験"),
            ("Day 6-10", "Gemini基本操作", "登録・実運用開始"),
            ("Day 11-15", "業務別AI活用", "メール・提案文・データ整理"),
            ("Day 16-20", "自動化・効率化", "プロンプト最適化"),
            ("Day 21-25", "スキルアップ", "高度な応用技法"),
            ("Day 26-30", "継続改善", "月間評価と改善計画"),
        ]

        y_pos = 1.3
        for i, (day, phase, content) in enumerate(roadmap):
            box = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.85))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(220, 245, 220) if i % 2 == 0 else self.colors['white']
            box.line.color.rgb = RGBColor(76, 175, 80)
            box.line.width = Pt(1)

            text_frame = box.text_frame
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_top = Inches(0.1)

            p = text_frame.paragraphs[0]
            p.text = f"{day}  |  {phase}  |  {content}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['text']

            y_pos += 1.0

        return prs

    def run(self):
        """実行"""
        print("[Claude Design] PowerPointプレゼン資料を生成中...\n")

        products = [
            ("AI時代の個人スキル販売術", self.create_sales_presentation),
            ("SNS運用自動化キット", self.create_sns_presentation),
            ("初心者向けAI活用ガイド", self.create_ai_guide_presentation)
        ]

        for product_name, generator_func in products:
            print(f"【{product_name}】")

            product_dir = Path(os.path.join(self.base_path, f"生成物・商品/素材/{product_name}"))
            pptx_file = product_dir / f"{product_name}_プレゼン資料.pptx"

            # PowerPointを生成
            prs = generator_func()

            # ファイルを保存
            prs.save(str(pptx_file))

            print(f"  ✅ {product_name}_プレゼン資料.pptx")
            print(f"     ファイルサイズ: {os.path.getsize(pptx_file) / 1024:.1f}KB")

        print("\n✅ PowerPointプレゼン資料生成完了")
        print("\n【次のステップ】")
        print("1. 各PowerPointファイルを確認")
        print("2. Excel自動化テンプレート作成")
        print("3. Gumroadにアップロード")

if __name__ == "__main__":
    designer = PresentationDesigner()
    designer.run()
